#!/usr/bin/env python3
"""Generate the A2A Protocol Demo presentation from SPEAKER_NOTES.md."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
BG_SLIDE     = RGBColor(0x13, 0x1C, 0x33)   # Slide background
ACCENT_BLUE  = RGBColor(0x42, 0x85, 0xF4)   # Google blue
ACCENT_GREEN = RGBColor(0x34, 0xA8, 0x53)   # Google green
ACCENT_YELLOW= RGBColor(0xFB, 0xBC, 0x04)   # Google yellow
ACCENT_RED   = RGBColor(0xEA, 0x43, 0x35)   # Google red
ACCENT_CYAN  = RGBColor(0x00, 0xBC, 0xD4)   # Cyan accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xBE, 0xC5)
MID_GRAY     = RGBColor(0x78, 0x90, 0x9C)
CODE_BG      = RGBColor(0x1A, 0x23, 0x3B)
ORANGE       = RGBColor(0xFF, 0x98, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ─────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_SLIDE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), space_after=Pt(2), font_name="Calibri", alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_bullet(tf, text, size=16, color=LIGHT_GRAY, level=0, bold=False, font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.level = level
    p.space_before = Pt(4)
    p.space_after = Pt(2)
    return p


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = add_rounded_rect(slide, left, top, width, height, CODE_BG, border_color=MID_GRAY)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    set_text(tf, code_text, size=font_size, color=ACCENT_CYAN, font_name="Consolas")
    return shape


def make_title_slide(title, subtitle="", section_num=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, BG_DARK)
    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)
    # Section number badge
    if section_num:
        badge = add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(1.2), Inches(0.55), ACCENT_BLUE)
        tf = badge.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_text(tf, section_num, size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    tb = add_text_box(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.5))
    set_text(tb.text_frame, title, size=40, color=WHITE, bold=True)
    # Subtitle
    if subtitle:
        tb2 = add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.0))
        set_text(tb2.text_frame, subtitle, size=20, color=LIGHT_GRAY)
    # Bottom bar
    add_shape_rect(slide, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), ACCENT_GREEN)
    return slide


def make_content_slide(title, section_label=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # Top accent bar
    add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.04), ACCENT_BLUE)
    # Title bar background
    add_shape_rect(slide, Inches(0), Inches(0.04), prs.slide_width, Inches(0.9), RGBColor(0x10, 0x19, 0x2D))
    # Section label (left)
    if section_label:
        tb = add_text_box(slide, Inches(0.6), Inches(0.15), Inches(2.0), Inches(0.6))
        set_text(tb.text_frame, section_label, size=11, color=ACCENT_BLUE, bold=True)
    # Title
    tb = add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.6))
    set_text(tb.text_frame, title, size=28, color=WHITE, bold=True)
    # Bottom bar
    add_shape_rect(slide, Inches(0), Inches(7.46), prs.slide_width, Inches(0.04), ACCENT_BLUE)
    return slide


def add_table(slide, left, top, width, height, rows, cols, data, header_color=ACCENT_BLUE):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = header_color
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x3B) if i % 2 == 1 else BG_SLIDE
    return table_shape


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

# ── SLIDE 1: Title ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)
add_shape_rect(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT_GREEN)

tb = add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11.0), Inches(1.2))
set_text(tb.text_frame, "A2A Protocol Demo", size=52, color=WHITE, bold=True)

tb = add_text_box(slide, Inches(1.0), Inches(2.8), Inches(11.0), Inches(0.8))
set_text(tb.text_frame, "A Complete Implementation of the Agent-to-Agent Protocol", size=26, color=ACCENT_CYAN)

tb = add_text_box(slide, Inches(1.0), Inches(4.0), Inches(11.0), Inches(1.5))
tf = tb.text_frame
set_text(tf, "Built with Google ADK 1.25.1  |  Vertex AI  |  Cloud Run", size=18, color=LIGHT_GRAY)
add_paragraph(tf, "24 Features  |  9 Agents  |  4 Auth Schemes  |  300 Tests", size=18, color=LIGHT_GRAY, space_before=Pt(12))

tb = add_text_box(slide, Inches(1.0), Inches(5.8), Inches(11.0), Inches(0.5))
set_text(tb.text_frame, "GCP Project: openshift-research-479702  |  Region: us-central1", size=14, color=MID_GRAY)


# ── SLIDE 2: Agenda ──────────────────────────────────────────────────────────
slide = make_content_slide("Agenda")
tb = add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.8))
tf = tb.text_frame
tf.word_wrap = True
items_left = [
    ("01", "What is A2A?", "The protocol & why it matters"),
    ("02", "Google ADK", "Agent Development Kit fundamentals"),
    ("03", "Architecture", "System design & communication flow"),
    ("04", "Project Structure", "Directory layout & shared infra"),
    ("05", "Agent Deep Dives", "All 9 agents explained"),
]
for num, title, desc in items_left:
    p = set_text(tf, "", size=14, color=WHITE) if tf.paragraphs[0].text == "" else add_paragraph(tf, "", size=14, color=WHITE, space_before=Pt(14))
    run = p.add_run()
    run.text = f"{num}  "
    run.font.size = Pt(14)
    run.font.color.rgb = ACCENT_BLUE
    run.font.bold = True
    run.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = title
    run2.font.size = Pt(18)
    run2.font.color.rgb = WHITE
    run2.font.bold = True
    run2.font.name = "Calibri"
    add_bullet(tf, desc, size=14, color=LIGHT_GRAY, level=0)

tb2 = add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.8))
tf2 = tb2.text_frame
tf2.word_wrap = True
items_right = [
    ("06", "Authentication", "4 schemes from none to OAuth 2.0"),
    ("07", "Callbacks & Safety", "Guardrails, logging, caching"),
    ("08", "Testing", "300 tests, mocking strategies"),
    ("09", "Deployment", "Local, Cloud Run, Agent Engine"),
    ("10", "Live Demo", "End-to-end walkthrough"),
]
for num, title, desc in items_right:
    p = set_text(tf2, "", size=14, color=WHITE) if tf2.paragraphs[0].text == "" else add_paragraph(tf2, "", size=14, color=WHITE, space_before=Pt(14))
    run = p.add_run()
    run.text = f"{num}  "
    run.font.size = Pt(14)
    run.font.color.rgb = ACCENT_GREEN
    run.font.bold = True
    run.font.name = "Calibri"
    run2 = p.add_run()
    run2.text = title
    run2.font.size = Pt(18)
    run2.font.color.rgb = WHITE
    run2.font.bold = True
    run2.font.name = "Calibri"
    add_bullet(tf2, desc, size=14, color=LIGHT_GRAY, level=0)


# ── SLIDE 3: The Problem ─────────────────────────────────────────────────────
make_title_slide(
    "Why A2A?",
    "AI agents today are siloed. A LangChain agent can't talk to a CrewAI agent.\n"
    "An AutoGen agent can't discover a Vertex AI agent's capabilities.\n"
    "Every framework has its own proprietary communication format.",
    "SECTION 01"
)

# ── SLIDE 4: A2A Protocol Core Concepts ──────────────────────────────────────
slide = make_content_slide("A2A Protocol  --  Core Concepts", "SECTION 01")

# 4 concept boxes in a 2x2 grid
concepts = [
    ("Agent Card", "/.well-known/agent.json", "JSON document advertising name,\nskills, capabilities, auth requirements.\nThink: business card for AI agents.", ACCENT_BLUE),
    ("JSON-RPC 2.0", "POST /", "All communication uses JSON-RPC 2.0.\nMethods: message/send, message/stream,\ntasks/get, tasks/cancel, tasks/list", ACCENT_GREEN),
    ("Task Lifecycle", "7 States", "submitted -> working -> completed\n             -> failed / canceled\n             -> input-required", ACCENT_YELLOW),
    ("Message Parts", "role + parts[]", "TextPart: plain text\nFilePart: binary data or URI\nDataPart: structured JSON", ORANGE),
]
positions = [
    (Inches(0.6), Inches(1.3)), (Inches(6.8), Inches(1.3)),
    (Inches(0.6), Inches(4.3)), (Inches(6.8), Inches(4.3)),
]
for (title, badge_text, desc, color), (left, top) in zip(concepts, positions):
    box = add_rounded_rect(slide, left, top, Inches(5.8), Inches(2.7), RGBColor(0x1A, 0x23, 0x3B), border_color=color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_top = Inches(0.2)
    set_text(tf, title, size=22, color=color, bold=True)
    add_paragraph(tf, badge_text, size=13, color=MID_GRAY, bold=True, font_name="Consolas", space_before=Pt(4))
    add_paragraph(tf, desc, size=14, color=LIGHT_GRAY, space_before=Pt(10))


# ── SLIDE 5: Protocol vs SDK vs Framework ────────────────────────────────────
slide = make_content_slide("Protocol  vs  SDK  vs  Framework", "SECTION 01")
data = [
    ["Layer", "What It Is", "Example"],
    ["A2A Protocol", "The specification (like HTTP spec)", "JSON-RPC methods, Agent Card schema"],
    ["a2a-sdk", "Reference implementation", "Python package, gRPC stubs"],
    ["Google ADK", "Agent framework that implements A2A", "LlmAgent, RemoteA2aAgent, to_a2a()"],
]
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.5), 4, 3, data)

tb = add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Key Insight", size=22, color=ACCENT_YELLOW, bold=True)
add_paragraph(tf, "You can build an A2A-compliant agent WITHOUT ADK.", size=18, color=WHITE, space_before=Pt(12))
add_paragraph(tf, "The a2a_client/client.py in this project proves it -- it uses only httpx (raw HTTP).", size=16, color=LIGHT_GRAY, space_before=Pt(8))
add_paragraph(tf, "The a2a_client/grpc_client.py does the same over gRPC with Protobuf.", size=16, color=LIGHT_GRAY, space_before=Pt(6))


# ── SLIDE 6: Google ADK Concepts ─────────────────────────────────────────────
make_title_slide(
    "Google Agent Development Kit (ADK)",
    "Open-source framework for building AI agents  |  Version 1.25.1",
    "SECTION 02"
)

# ── SLIDE 7: ADK Building Blocks ─────────────────────────────────────────────
slide = make_content_slide("ADK Building Blocks", "SECTION 02")

blocks = [
    ("LlmAgent", "LLM-backed agent with model, instruction,\ntools, sub_agents, callbacks, output_key", ACCENT_BLUE),
    ("SequentialAgent", "Runs sub-agents in order: A -> B -> C\nDeterministic pipeline orchestration", ACCENT_GREEN),
    ("ParallelAgent", "Runs sub-agents concurrently (fan-out)\nAll execute simultaneously", ACCENT_YELLOW),
    ("LoopAgent", "Repeats sub-agents until exit condition\nor max_iterations reached", ORANGE),
    ("RemoteA2aAgent", "Proxy to a remote A2A server. Discovers\ncapabilities via Agent Card at construction", ACCENT_CYAN),
    ("to_a2a()", "Converts any ADK agent into an A2A-compliant\nweb app with JSON-RPC + agent card route", ACCENT_RED),
]
for i, (name, desc, color) in enumerate(blocks):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.2)
    top = Inches(1.3 + row * 3.0)
    box = add_rounded_rect(slide, left, top, Inches(3.9), Inches(2.6), RGBColor(0x1A, 0x23, 0x3B), border_color=color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    set_text(tf, name, size=20, color=color, bold=True, font_name="Consolas")
    add_paragraph(tf, desc, size=14, color=LIGHT_GRAY, space_before=Pt(10))


# ── SLIDE 8: Architecture Overview ───────────────────────────────────────────
make_title_slide(
    "Architecture Overview",
    "Multi-agent ecosystem with specialist routing",
    "SECTION 03"
)

# ── SLIDE 9: Architecture Diagram (Top-Down Request Flow) ────────────────────
slide = make_content_slide("Request Flow  --  Top to Bottom", "SECTION 03")
code = (
    '  "What is the weather in Paris?"\n'
    "                |\n"
    "  Layer 1       v\n"
    "  CLIENT:   a2a_client (HTTP / gRPC)          -- Any HTTP client, no SDK needed\n"
    "                |\n"
    "                | JSON-RPC 2.0 over HTTP\n"
    "  Layer 2       v\n"
    "  ROUTER:   orchestrator_agent                 -- LLM reads descriptions, picks agent\n"
    "            5x RemoteA2aAgent sub-agents\n"
    "            httpx clients with auth headers\n"
    "                |\n"
    "                | A2A Protocol (Agent Card discovery + JSON-RPC)\n"
    "  Layer 3       v\n"
    "  AGENTS:   weather  research   code     data     async\n"
    "            :8001    :8002      :8003    :8004    :8005\n"
    "            No Auth  JWT       API Key  OAuth    Push+SSE\n"
    "                |\n"
    "  Layer 4       v\n"
    "  SHARED:   config.py  |  auth.py  |  callbacks.py\n"
    "            Settings     4 schemes    Logging + Guardrails + Cache"
)
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.8), code, font_size=14)


# ── SLIDE 9b: Full Architecture Diagram (PNG) ────────────────────────────────
slide = make_content_slide("Architecture Diagram", "SECTION 03")
import os
_diagram_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "architecture_diagram.png")
if os.path.exists(_diagram_path):
    slide.shapes.add_picture(_diagram_path, Inches(0.3), Inches(1.05), width=Inches(12.7))

# ── SLIDE 10: Communication Patterns ─────────────────────────────────────────
slide = make_content_slide("Communication Patterns", "SECTION 03")
data = [
    ["Pattern", "Example", "A2A Method", "Agent"],
    ["Sync request/response", '"Weather in Paris?"', "message/send", "weather_agent"],
    ["SSE Streaming", '"5-day forecast for Tokyo"', "message/stream", "weather, research"],
    ["Async + Push", '"Run 20-sec simulation"', "message/send + webhook", "async_agent"],
    ["Polling", "Loop agent checks status", "tasks/get", "loop_agent"],
    ["Cancellation", "Stop a running task", "tasks/cancel", "async_agent"],
    ["Multi-turn", '"Research AI" -> "Focus safety"', "message/send + taskId", "research_agent"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.0), 7, 4, data)


# ── SLIDE 11: Project Structure ──────────────────────────────────────────────
make_title_slide(
    "Project Structure",
    "Directory layout, shared infrastructure, and configuration",
    "SECTION 04"
)

# ── SLIDE 12: Directory Layout ───────────────────────────────────────────────
slide = make_content_slide("Directory Layout", "SECTION 04")
code_left = (
    "A2A_Research_demo_1/\n"
    "|\n"
    "|-- shared/              Config, auth, callbacks\n"
    "|-- weather_agent/       Simplest A2A agent\n"
    "|-- research_agent/      JWT + extended card\n"
    "|-- code_agent/          Code execution sandbox\n"
    "|-- data_agent/          CSV artifacts + OAuth\n"
    "|-- async_agent/         Custom FastAPI (no ADK)\n"
    "|-- orchestrator_agent/  Root router agent\n"
    "|-- pipeline_agent/      Sequential pipeline\n"
    "|-- parallel_agent/      Fan-out weather\n"
    "|-- loop_agent/          Polling loop"
)
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.5), code_left, font_size=13)

code_right = (
    "|-- a2a_client/          Standalone clients\n"
    "|   |-- client.py        HTTP (httpx only)\n"
    "|   |-- grpc_client.py   gRPC (protobuf)\n"
    "|\n"
    "|-- webhook_server/      Push notification sink\n"
    "|-- evals/               ADK evaluation datasets\n"
    "|-- protos/              gRPC .proto definitions\n"
    "|-- scripts/             start/stop/deploy\n"
    "|-- tests/               300 pytest tests\n"
    "|\n"
    "|-- .env                 Runtime configuration\n"
    "|-- requirements.txt     Pinned dependencies"
)
add_code_block(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), code_right, font_size=13)


# ── SLIDE 13: Shared Config ──────────────────────────────────────────────────
slide = make_content_slide("shared/config.py  --  Centralized Configuration", "SECTION 04")
code = (
    '@dataclass\n'
    'class Settings:\n'
    '    GOOGLE_CLOUD_PROJECT: str = field(\n'
    '        default_factory=lambda: os.environ.get("GOOGLE_CLOUD_PROJECT", "")\n'
    '    )\n'
    '    GEMINI_MODEL: str = "gemini-2.0-flash"\n'
    '    WEATHER_AGENT_URL: str = "http://localhost:8001"\n'
    '    # ... 14 fields total\n'
    '\n'
    '    def validate(self) -> None:\n'
    '        """Check all required env vars are set."""\n'
    '        ...\n'
    '\n'
    'settings = Settings()            # Singleton\n'
    'if "pytest" not in sys.modules:  # Skip during test collection\n'
    '    settings.validate()'
)
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(7.5), Inches(5.0), code, font_size=14)

tb = add_text_box(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Design Principles", size=20, color=ACCENT_BLUE, bold=True)
add_bullet(tf, "Python dataclass -- typed, testable", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Reads from .env via python-dotenv", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Singleton pattern -- one instance", size=14, color=LIGHT_GRAY)
add_bullet(tf, "validate() checks required fields", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Test guard prevents broken collection", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "Every agent imports:", size=14, color=ACCENT_YELLOW, bold=True, space_before=Pt(20))
add_paragraph(tf, "from shared.config import settings", size=13, color=ACCENT_CYAN, font_name="Consolas", space_before=Pt(6))


# ── SLIDE 14: weather_agent ──────────────────────────────────────────────────
make_title_slide(
    "Agent Deep Dives",
    "Starting with the simplest: weather_agent",
    "SECTION 05"
)

# ── SLIDE 15: weather_agent code ─────────────────────────────────────────────
slide = make_content_slide("weather_agent  --  The Minimal A2A Agent Pattern", "SECTION 05")
code = (
    '# Step 1: Define the Agent Card\n'
    '_AGENT_CARD = AgentCard(\n'
    '    name="weather_agent",\n'
    '    skills=[_weather_skill, _forecast_skill],\n'
    '    capabilities=AgentCapabilities(streaming=True),\n'
    ')\n'
    '\n'
    '# Step 2: Create the LLM Agent\n'
    'root_agent = LlmAgent(\n'
    '    model=settings.GEMINI_MODEL,\n'
    '    instruction="You are a weather assistant...",\n'
    '    tools=[get_weather, get_forecast],\n'
    ')\n'
    '\n'
    '# Step 3: Convert to A2A web application\n'
    'app = to_a2a(root_agent, port=8001, agent_card=_AGENT_CARD)\n'
    '\n'
    '# That\'s it! to_a2a() auto-creates:\n'
    '#   GET  /.well-known/agent.json  --> Agent Card\n'
    '#   POST /                        --> JSON-RPC dispatcher'
)
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(7.5), Inches(5.5), code, font_size=14)

tb = add_text_box(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "3-Step Pattern", size=22, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "1. Agent Card", size=18, color=WHITE, bold=True, space_before=Pt(16))
add_bullet(tf, "Advertise skills & capabilities", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "2. LLM Agent", size=18, color=WHITE, bold=True, space_before=Pt(12))
add_bullet(tf, "Model + instruction + tools", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "3. to_a2a()", size=18, color=WHITE, bold=True, space_before=Pt(12))
add_bullet(tf, "Instant A2A web server", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "Features: F1, F2, F3, F8, F12", size=14, color=ACCENT_YELLOW, space_before=Pt(20))


# ── SLIDE 16: Function Tools ─────────────────────────────────────────────────
slide = make_content_slide("Function Tools  --  How ADK Turns Python into Tool Calls", "SECTION 05")
code = (
    'async def get_weather(city: str) -> dict:\n'
    '    """\n'
    '    Return current weather conditions for a city.\n'
    '    Args:\n'
    '        city: City name, e.g. "London"\n'
    '    Returns:\n'
    '        Dict with temperature_c, conditions, etc.\n'
    '    """\n'
    '    if not settings.OPENWEATHERMAP_API_KEY:\n'
    '        return _mock_weather(city)  # Fallback\n'
    '\n'
    '    async with httpx.AsyncClient() as client:\n'
    '        resp = await client.get(f"{OWM}/weather",\n'
    '                                params={"q": city, ...})\n'
    '        data = resp.json()\n'
    '        return {"city": data["name"],\n'
    '                "temperature_c": data["main"]["temp"], ...}'
)
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(7.5), Inches(5.5), code, font_size=14)

tb = add_text_box(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "ADK Auto-Schema", size=20, color=ACCENT_BLUE, bold=True)
add_bullet(tf, "Type hints -> JSON Schema", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Docstring -> tool description", size=14, color=LIGHT_GRAY)
add_bullet(tf, "LLM sees function signature", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "Mock Data Fallback", size=20, color=ACCENT_YELLOW, bold=True, space_before=Pt(20))
add_bullet(tf, "No API key? Return mock data", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Demo works out of the box", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "Async by Design", size=20, color=ACCENT_GREEN, bold=True, space_before=Pt(20))
add_bullet(tf, "httpx.AsyncClient -- non-blocking", size=14, color=LIGHT_GRAY)
add_bullet(tf, "ADK handles async tools natively", size=14, color=LIGHT_GRAY)


# ── SLIDE 17: orchestrator_agent ─────────────────────────────────────────────
slide = make_content_slide("orchestrator_agent  --  The Router", "SECTION 05")
code = (
    '# Pre-configured httpx clients with auth headers\n'
    '_code_client = httpx.AsyncClient(\n'
    '    headers={"X-API-Key": settings.CODE_AGENT_API_KEY},\n'
    '    timeout=httpx.Timeout(120.0),\n'
    ')\n'
    '\n'
    '# RemoteA2aAgent with auth client\n'
    'code_agent = RemoteA2aAgent(\n'
    '    name="code_agent",\n'
    '    agent_card=f"{URL}/.well-known/agent.json",\n'
    '    httpx_client=_code_client,  # auth pre-set\n'
    ')\n'
    '\n'
    'root_agent = LlmAgent(\n'
    '    model=settings.GEMINI_MODEL,\n'
    '    instruction="Route to the right specialist...",\n'
    '    sub_agents=[weather, research, code, data, async],\n'
    '    tools=[list_available_agents, get_agent_status],\n'
    ')'
)
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(7.5), Inches(4.8), code, font_size=14)

tb = add_text_box(slide, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "LLM-as-Router", size=20, color=ACCENT_BLUE, bold=True)
add_bullet(tf, "LLM reads sub-agent descriptions", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Decides which specialist to call", size=14, color=LIGHT_GRAY)
add_bullet(tf, "No if/else routing rules", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "Auth via httpx Clients", size=20, color=ACCENT_GREEN, bold=True, space_before=Pt(16))
add_bullet(tf, "Each RemoteA2aAgent gets a", size=14, color=LIGHT_GRAY)
add_bullet(tf, "pre-configured httpx.AsyncClient", size=14, color=LIGHT_GRAY)
add_bullet(tf, "with auth headers + 120s timeout", size=14, color=LIGHT_GRAY)


# ── SLIDE 18: async_agent ────────────────────────────────────────────────────
slide = make_content_slide("async_agent  --  Custom A2A (No ADK)", "SECTION 05")

tb = add_text_box(slide, Inches(0.6), Inches(1.3), Inches(6.0), Inches(5.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Why Custom?", size=20, color=ACCENT_YELLOW, bold=True)
add_paragraph(tf, "ADK's default handler doesn't support:", size=15, color=LIGHT_GRAY, space_before=Pt(8))
add_bullet(tf, "Background task execution (20-sec tasks)", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Push notification delivery (webhooks)", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Task cancellation (asyncio.Task.cancel())", size=14, color=LIGHT_GRAY)
add_bullet(tf, "SSE streaming with progress updates", size=14, color=LIGHT_GRAY)
add_bullet(tf, "Cursor-based task pagination", size=14, color=LIGHT_GRAY)

add_paragraph(tf, "In-Memory Architecture", size=20, color=ACCENT_BLUE, bold=True, space_before=Pt(16))
add_bullet(tf, "_task_store: dict[str, dict]  -- task state", size=13, color=ACCENT_CYAN, font_name="Consolas")
add_bullet(tf, "_webhook_store: dict[str, dict]  -- configs", size=13, color=ACCENT_CYAN, font_name="Consolas")
add_bullet(tf, "_running_tasks: dict[str, asyncio.Task]", size=13, color=ACCENT_CYAN, font_name="Consolas")
add_bullet(tf, "_sse_queues: dict[str, list[Queue]]", size=13, color=ACCENT_CYAN, font_name="Consolas")

tb2 = add_text_box(slide, Inches(6.9), Inches(1.3), Inches(5.9), Inches(5.8))
tf2 = tb2.text_frame
tf2.word_wrap = True
set_text(tf2, "Request Flow", size=20, color=ACCENT_GREEN, bold=True)
flow_items = [
    "1. Client POSTs message/send",
    "2. Create task (status: submitted)",
    "3. asyncio.create_task() for background",
    "4. Store Task in _running_tasks",
    "5. Return immediately to client",
    "6. Background: 20-sec simulation",
    "   - 0%, 25%, 50%, 75% progress",
    "   - Push notifications to webhook",
    "   - SSE events to connected clients",
    "7. Final: status -> completed + artifact",
]
for item in flow_items:
    add_bullet(tf2, item, size=14, color=LIGHT_GRAY)


# ── SLIDE 19: Workflow Agents ─────────────────────────────────────────────────
slide = make_content_slide("Workflow Agents  --  Sequential, Parallel, Loop", "SECTION 05")

# Pipeline
box = add_rounded_rect(slide, Inches(0.6), Inches(1.3), Inches(3.8), Inches(5.5), RGBColor(0x1A, 0x23, 0x3B), border_color=ACCENT_BLUE)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
set_text(tf, "pipeline_agent", size=18, color=ACCENT_BLUE, bold=True, font_name="Consolas")
add_paragraph(tf, "SequentialAgent", size=14, color=MID_GRAY, space_before=Pt(4))
add_paragraph(tf, "fetch_agent", size=16, color=WHITE, bold=True, space_before=Pt(16))
add_bullet(tf, "Gather raw information", size=13, color=LIGHT_GRAY)
add_bullet(tf, "output_key='raw_data'", size=12, color=ACCENT_CYAN, font_name="Consolas")
add_paragraph(tf, "        v", size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(4))
add_paragraph(tf, "analyze_agent", size=16, color=WHITE, bold=True, space_before=Pt(4))
add_bullet(tf, "Extract 3-5 key insights", size=13, color=LIGHT_GRAY)
add_bullet(tf, "output_key='analysis'", size=12, color=ACCENT_CYAN, font_name="Consolas")
add_paragraph(tf, "        v", size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(4))
add_paragraph(tf, "report_agent", size=16, color=WHITE, bold=True, space_before=Pt(4))
add_bullet(tf, "Markdown report", size=13, color=LIGHT_GRAY)
add_bullet(tf, "output_key='final_report'", size=12, color=ACCENT_CYAN, font_name="Consolas")

# Parallel
box = add_rounded_rect(slide, Inches(4.7), Inches(1.3), Inches(3.8), Inches(5.5), RGBColor(0x1A, 0x23, 0x3B), border_color=ACCENT_GREEN)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
set_text(tf, "parallel_agent", size=18, color=ACCENT_GREEN, bold=True, font_name="Consolas")
add_paragraph(tf, "ParallelAgent + SequentialAgent", size=14, color=MID_GRAY, space_before=Pt(4))
add_paragraph(tf, "Fan-Out (concurrent):", size=16, color=WHITE, bold=True, space_before=Pt(14))
add_bullet(tf, "London   -> RemoteA2aAgent", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Tokyo    -> RemoteA2aAgent", size=13, color=LIGHT_GRAY)
add_bullet(tf, "New York -> RemoteA2aAgent", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Sydney   -> RemoteA2aAgent", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Paris    -> RemoteA2aAgent", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "        v", size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(4))
add_paragraph(tf, "Fan-In (aggregate):", size=16, color=WHITE, bold=True, space_before=Pt(4))
add_bullet(tf, "Summary table of all 5 cities", size=13, color=LIGHT_GRAY)

# Loop
box = add_rounded_rect(slide, Inches(8.8), Inches(1.3), Inches(3.9), Inches(5.5), RGBColor(0x1A, 0x23, 0x3B), border_color=ACCENT_YELLOW)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
set_text(tf, "loop_agent", size=18, color=ACCENT_YELLOW, bold=True, font_name="Consolas")
add_paragraph(tf, "LoopAgent (max 10 iterations)", size=14, color=MID_GRAY, space_before=Pt(4))
add_paragraph(tf, "start_task_agent", size=16, color=WHITE, bold=True, space_before=Pt(14))
add_bullet(tf, "Delegates to RemoteA2aAgent", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Starts async task", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "        v", size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(4))
add_paragraph(tf, "poll_agent (loop)", size=16, color=WHITE, bold=True, space_before=Pt(4))
add_bullet(tf, "Checks task status", size=13, color=LIGHT_GRAY)
add_bullet(tf, "WORKING -> continue", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "exit_check_agent", size=16, color=WHITE, bold=True, space_before=Pt(8))
add_bullet(tf, "DONE/FAILED -> EXIT", size=13, color=LIGHT_GRAY)
add_bullet(tf, "WORKING -> CONTINUE", size=13, color=LIGHT_GRAY)


# ── SLIDE 20: Auth Schemes ───────────────────────────────────────────────────
make_title_slide(
    "Authentication Schemes",
    "4 schemes from zero security to full OAuth 2.0",
    "SECTION 06"
)

# ── SLIDE 21: Auth Comparison ─────────────────────────────────────────────────
slide = make_content_slide("Four Authentication Schemes", "SECTION 06")
data = [
    ["Scheme", "Agent", "Header", "How It Works"],
    ["No Auth", "weather_agent", "(none)", "Fully open -- no credentials needed"],
    ["API Key", "code_agent", "X-API-Key: <key>", "Shared secret in header, middleware validates"],
    ["Bearer JWT", "research_agent", "Authorization: Bearer <jwt>", "HMAC-SHA256 signed token with expiry"],
    ["OAuth 2.0", "data_agent", "Authorization: Bearer <token>", "GCP service account token verification"],
]
add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.0), 5, 4, data)

tb = add_text_box(slide, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Cross-Cutting Rule", size=20, color=ACCENT_YELLOW, bold=True)
add_paragraph(tf, "Discovery is ALWAYS public. Operations require auth.", size=18, color=WHITE, space_before=Pt(10))
add_paragraph(tf, 'if request.url.path == "/.well-known/agent.json":  return await call_next(request)', size=13, color=ACCENT_CYAN, font_name="Consolas", space_before=Pt(10))
add_paragraph(tf, "Agents must be discoverable without pre-shared credentials -- this is essential for A2A.", size=14, color=LIGHT_GRAY, space_before=Pt(8))


# ── SLIDE 22: Callbacks ──────────────────────────────────────────────────────
make_title_slide(
    "Callbacks, Guardrails & Safety",
    "6 ADK callback hooks for logging, security, and caching",
    "SECTION 07"
)

# ── SLIDE 23: Callback Chain ─────────────────────────────────────────────────
slide = make_content_slide("The Callback Chain", "SECTION 07")
code = (
    "Incoming message\n"
    "     |\n"
    "     v\n"
    "[before_model_callback]  <-- Log agent + message count\n"
    "     |\n"
    "     v\n"
    "  LLM Call (Gemini)\n"
    "     |\n"
    "     v\n"
    "[after_model_callback]   <-- Log tokens, redact URLs\n"
    "     |\n"
    "     v\n"
    "  LLM decides to call a tool\n"
    "     |\n"
    "     v\n"
    "[before_tool_callback]   <-- GUARDRAIL: block os.system, eval\n"
    "     |                       CACHE: return cached result\n"
    "     v\n"
    "  Tool Execution\n"
    "     |\n"
    "     v\n"
    "[after_tool_callback]    <-- Log result, store in cache\n"
    "     |\n"
    "     v\n"
    "Response to client"
)
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(7.0), Inches(5.8), code, font_size=13)

tb = add_text_box(slide, Inches(8.0), Inches(1.5), Inches(4.8), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Return Value Semantics", size=20, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "return None", size=16, color=ACCENT_GREEN, bold=True, font_name="Consolas", space_before=Pt(14))
add_bullet(tf, "Pass through (normal processing)", size=14, color=LIGHT_GRAY)
add_paragraph(tf, "return {...}", size=16, color=ACCENT_RED, bold=True, font_name="Consolas", space_before=Pt(12))
add_bullet(tf, "Intercept (use this result instead)", size=14, color=LIGHT_GRAY)

add_paragraph(tf, "Guardrail Example", size=20, color=ACCENT_RED, bold=True, space_before=Pt(20))
add_paragraph(tf, "Blocked patterns:", size=14, color=LIGHT_GRAY, space_before=Pt(6))
add_bullet(tf, "os.system, subprocess", size=13, color=ACCENT_CYAN, font_name="Consolas")
add_bullet(tf, "eval(), exec(), open()", size=13, color=ACCENT_CYAN, font_name="Consolas")
add_bullet(tf, "shutil.rmtree, __import__", size=13, color=ACCENT_CYAN, font_name="Consolas")


# ── SLIDE 24: Testing ────────────────────────────────────────────────────────
make_title_slide(
    "Testing Strategy",
    "300 tests  |  ~4 seconds  |  Zero GCP calls",
    "SECTION 08"
)

# ── SLIDE 25: Testing Details ─────────────────────────────────────────────────
slide = make_content_slide("Test Architecture", "SECTION 08")
data = [
    ["Category", "Files", "What They Test"],
    ["Config", "test_config.py", "Settings defaults, validation, env overrides"],
    ["Auth", "test_shared_auth.py", "API key, Bearer JWT, HMAC signatures"],
    ["Callbacks", "test_shared_callbacks.py", "Logging, guardrails, caching"],
    ["Weather", "test_weather_agent.py + ext", "Mock data, API parsing, aggregation"],
    ["Async", "test_async_agent.py + lifecycle", "Task CRUD, pagination, cancel, HMAC"],
    ["Webhook", "test_webhook_server.py + ext", "Event receipt, HMAC, timestamps"],
    ["Client", "test_a2a_client.py + ext", "URL normalization, send, stream, get"],
    ["Data", "test_data_agent.py + ext", "CSV parsing, statistics, reports"],
    ["Orchestrator", "test_orchestrator_*.py", "Agent listing, status, URL redaction"],
]
add_table(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), 10, 3, data)


# ── SLIDE 26: Testing Techniques ─────────────────────────────────────────────
slide = make_content_slide("Key Testing Techniques", "SECTION 08")
code = (
    '# 1. Auto-mock environment (conftest.py)\n'
    '@pytest.fixture(autouse=True)\n'
    'def mock_env_vars(monkeypatch):\n'
    '    monkeypatch.setenv("GOOGLE_GENAI_USE_VERTEXAI", "0")\n'
    '    monkeypatch.setenv("OPENWEATHERMAP_API_KEY", "")  # mock\n'
    '\n'
    '# 2. FastAPI TestClient (sync wrapper for async)\n'
    'client = TestClient(app)\n'
    'resp = client.post("/", json=payload)\n'
    'assert resp.status_code == 200\n'
    '\n'
    '# 3. AsyncMock for httpx.AsyncClient\n'
    'mock_client = AsyncMock()\n'
    'mock_client.get = AsyncMock(return_value=mock_resp)\n'
    'with patch("httpx.AsyncClient", return_value=mock_client):\n'
    '    result = await tools_mod.get_weather("London")\n'
    '\n'
    '# 4. monkeypatch settings singleton\n'
    'monkeypatch.setattr(tools_mod.settings, "API_KEY", "fake")'
)
add_code_block(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.5), code, font_size=14)


# ── SLIDE 27: Deployment ─────────────────────────────────────────────────────
make_title_slide(
    "Deployment",
    "Local dev  |  Google Cloud Run  |  Vertex AI Agent Engine",
    "SECTION 09"
)

# ── SLIDE 28: Deployment Details ──────────────────────────────────────────────
slide = make_content_slide("Deployment Options", "SECTION 09")

# Local
box = add_rounded_rect(slide, Inches(0.6), Inches(1.3), Inches(3.8), Inches(5.5), RGBColor(0x1A, 0x23, 0x3B), border_color=ACCENT_BLUE)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
set_text(tf, "Local Dev", size=20, color=ACCENT_BLUE, bold=True)
add_paragraph(tf, "./scripts/start_all.sh", size=13, color=ACCENT_CYAN, font_name="Consolas", space_before=Pt(12))
add_bullet(tf, "weather_agent  :8001", size=13, color=LIGHT_GRAY)
add_bullet(tf, "research_agent :8002", size=13, color=LIGHT_GRAY)
add_bullet(tf, "code_agent     :8003", size=13, color=LIGHT_GRAY)
add_bullet(tf, "data_agent     :8004", size=13, color=LIGHT_GRAY)
add_bullet(tf, "async_agent    :8005", size=13, color=LIGHT_GRAY)
add_bullet(tf, "webhook_server :9000", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "adk web --host 0.0.0.0 --port 8000 .", size=13, color=ACCENT_CYAN, font_name="Consolas", space_before=Pt(12))
add_bullet(tf, "ADK Dev UI at :8000/dev-ui/", size=13, color=LIGHT_GRAY)

# Cloud Run
box = add_rounded_rect(slide, Inches(4.7), Inches(1.3), Inches(3.8), Inches(5.5), RGBColor(0x1A, 0x23, 0x3B), border_color=ACCENT_GREEN)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
set_text(tf, "Cloud Run", size=20, color=ACCENT_GREEN, bold=True)
add_paragraph(tf, "Dockerised microservices", size=14, color=LIGHT_GRAY, space_before=Pt(12))
add_paragraph(tf, "./scripts/deploy_cloud_run.sh", size=13, color=ACCENT_CYAN, font_name="Consolas", space_before=Pt(10))
add_bullet(tf, "Builds with Cloud Build", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Auto-scales to zero", size=13, color=LIGHT_GRAY)
add_bullet(tf, "HTTPS endpoints", size=13, color=LIGHT_GRAY)
add_bullet(tf, "IAM authentication", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "Each agent has its own", size=14, color=LIGHT_GRAY, space_before=Pt(12))
add_paragraph(tf, "Dockerfile", size=14, color=ACCENT_YELLOW, bold=True, space_before=Pt(4))

# Agent Engine
box = add_rounded_rect(slide, Inches(8.8), Inches(1.3), Inches(3.9), Inches(5.5), RGBColor(0x1A, 0x23, 0x3B), border_color=ACCENT_YELLOW)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.2)
set_text(tf, "Vertex AI Agent Engine", size=20, color=ACCENT_YELLOW, bold=True)
add_paragraph(tf, "Managed agent hosting", size=14, color=LIGHT_GRAY, space_before=Pt(12))
add_bullet(tf, "Orchestrator deployed as", size=13, color=LIGHT_GRAY)
add_bullet(tf, "managed Vertex AI agent", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Auto-routing to Cloud Run", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Built-in monitoring", size=13, color=LIGHT_GRAY)
add_bullet(tf, "OpenTelemetry tracing", size=13, color=LIGHT_GRAY)
add_paragraph(tf, "GCP Resources:", size=14, color=ACCENT_BLUE, bold=True, space_before=Pt(14))
add_bullet(tf, "Cloud Run APIs", size=13, color=LIGHT_GRAY)
add_bullet(tf, "Cloud Build APIs", size=13, color=LIGHT_GRAY)
add_bullet(tf, "GCS staging bucket", size=13, color=LIGHT_GRAY)


# ── SLIDE 29: 24 Features Matrix ─────────────────────────────────────────────
slide = make_content_slide("All 24 A2A Features", "COMPLETE MATRIX")
data = [
    ["#", "Feature", "Where"],
    ["F1", "Agent Cards", "All agents"],
    ["F2", "Sync Request/Response", "All agents"],
    ["F3", "SSE Streaming", "weather, research, async"],
    ["F4", "Push Notifications", "async -> webhook_server"],
    ["F5", "Task Lifecycle", "async_agent (7 states)"],
    ["F6", "Multi-turn", "research, loop agents"],
    ["F7", "Extended Card", "research_agent"],
    ["F8", "Auth Schemes (4)", "weather/code/research/data"],
    ["F9", "A2A Routing", "orchestrator_agent"],
    ["F10", "Workflow Agents", "pipeline/parallel/loop"],
    ["F11", "Agent Types (5)", "LLM/Seq/Para/Loop/Remote"],
    ["F12", "Tool Types (3)", "Functions/Search/CodeExec"],
]
add_table(slide, Inches(0.3), Inches(1.2), Inches(6.3), Inches(5.8), 13, 3, data)

data2 = [
    ["#", "Feature", "Where"],
    ["F13", "Session State", "pipeline_agent"],
    ["F14", "Memory", "research_agent"],
    ["F15", "Artifacts", "data_agent (CSV)"],
    ["F16", "Callbacks (6 hooks)", "shared/callbacks.py"],
    ["F17", "Safety Guardrails", "code_agent"],
    ["F18", "Evaluation", "evals/ (5 datasets)"],
    ["F19", "Agent Engine", "orchestrator deploy"],
    ["F20", "Cloud Run", "All remote agents"],
    ["F21", "gRPC Transport", "a2a_client/grpc_client"],
    ["F22", "Observability", "OpenTelemetry config"],
    ["F23", "ADK Dev UI", "adk run/web/api_server"],
    ["F24", "Interoperability", "a2a_client (httpx only)"],
]
add_table(slide, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.8), 13, 3, data2)


# ── SLIDE 30: Live Demo ──────────────────────────────────────────────────────
make_title_slide(
    "Live Demo",
    "Agent discovery  |  Routing  |  Async tasks  |  Auth  |  Tests",
    "SECTION 10"
)

# ── SLIDE 31: Demo Commands ──────────────────────────────────────────────────
slide = make_content_slide("Demo Commands", "SECTION 10")
code = (
    '# 1. Start all agents (uvicorn, not adk api_server)\n'
    './scripts/start_all.sh\n'
    '\n'
    '# 2. Agent Discovery\n'
    'curl http://localhost:8001/.well-known/agent.json | python3 -m json.tool\n'
    '\n'
    '# 3. Sync message/send (note: messageId is required)\n'
    'curl -X POST http://localhost:8001/ \\\n'
    '  -H "Content-Type: application/json" \\\n'
    '  -d \'{"jsonrpc":"2.0","id":"1","method":"message/send",\n'
    '       "params":{"message":{"role":"user","messageId":"m1",\n'
    '       "parts":[{"kind":"text","text":"Weather in Paris?"}]}}}\'\n'
    '\n'
    '# 4. Orchestrator (ADK Web UI -- use "." for project root)\n'
    'adk web --host 0.0.0.0 --port 8000 .\n'
    '\n'
    '# 5. Standalone A2A client (no ADK, just httpx)\n'
    'python3 -m a2a_client.client\n'
    '\n'
    '# 6. Run Tests\n'
    'pytest tests/ -v'
)
add_code_block(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.8), code, font_size=14)


# ── SLIDE 32: Key Takeaways ──────────────────────────────────────────────────
slide = make_content_slide("Key Takeaways")

takeaways = [
    ("A2A is the HTTP of the agent world", "Open protocol for agent-to-agent communication, framework-agnostic", ACCENT_BLUE),
    ("Agent Cards enable discovery", "Standardized capability advertisement at /.well-known/agent.json", ACCENT_GREEN),
    ("ADK makes it simple", "3 steps: Agent Card -> LlmAgent -> to_a2a() = production A2A server", ACCENT_YELLOW),
    ("Security is layered", "Discovery always public, operations authenticated, 4 auth schemes", ACCENT_RED),
    ("Test everything", "300 tests, mock environment, zero GCP calls in CI", ACCENT_CYAN),
    ("Production-ready patterns", "Callbacks, guardrails, URL redaction, HMAC signatures, retry logic", ORANGE),
]
for i, (title, desc, color) in enumerate(takeaways):
    col = i % 2
    row = i // 2
    left = Inches(0.6 + col * 6.3)
    top = Inches(1.3 + row * 2.0)
    box = add_rounded_rect(slide, left, top, Inches(5.9), Inches(1.7), RGBColor(0x1A, 0x23, 0x3B), border_color=color)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    set_text(tf, title, size=18, color=color, bold=True)
    add_paragraph(tf, desc, size=13, color=LIGHT_GRAY, space_before=Pt(6))


# ── SLIDE 33: Thank You ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_BLUE)
add_shape_rect(slide, Inches(0), Inches(7.42), prs.slide_width, Inches(0.08), ACCENT_GREEN)

tb = add_text_box(slide, Inches(1.0), Inches(2.0), Inches(11.0), Inches(1.0))
set_text(tb.text_frame, "Thank You", size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tb = add_text_box(slide, Inches(1.0), Inches(3.5), Inches(11.0), Inches(1.5))
tf = tb.text_frame
set_text(tf, "Resources", size=24, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "A2A Protocol Spec:  a2a-protocol.org/latest/specification/", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(14))
add_paragraph(tf, "Google ADK Docs:  google.github.io/adk-docs/", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(8))
add_paragraph(tf, "Project Docs:  README.md  |  ENV_SETUP.md  |  DEMO.md  |  SPEAKER_NOTES.md", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(8))


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════

output_path = "/home/sandbox1/A2A_Research_demo_1/A2A_Protocol_Demo.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
